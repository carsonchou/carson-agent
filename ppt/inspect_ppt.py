# -*- coding: utf-8 -*-
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

SRC = r'C:\Users\User\Downloads\carson-agent\.playwright-mcp\期末專題PPT.pptx'
prs = Presentation(SRC)

print(f'投影片數量: {len(prs.slides)}')
print(f'尺寸: {prs.slide_width.inches:.2f}" x {prs.slide_height.inches:.2f}"')
print()

for i, slide in enumerate(prs.slides):
    print(f'=== S{i+1:02d} ===')
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                # 取前80字
                preview = text[:80].replace('\n', ' | ')
                print(f'  [{shape.shape_type}] {preview}')
    print()
