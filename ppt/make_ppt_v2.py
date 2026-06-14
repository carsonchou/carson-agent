# -*- coding: utf-8 -*-
"""
期末專題 PPT v2 — 現代化視覺設計
主題：以 UC3845 PWM IC 為核心之 BUCK 降壓轉換器電路研製
"""
import sys, io, fitz
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ─── 路徑 ──────────────────────────────────────────────────
DL  = r"C:\Users\User\Downloads"
OUT     = r"C:\Users\User\Desktop\期末專題PPT.pptx"
OUT_UPL = r"C:\Users\User\Downloads\carson-agent\.playwright-mcp\期末專題PPT_v2.pptx"
FT  = "標楷體"
FT2 = "Arial"          # 數字/英文用

# ─── 色票（Modern Navy × Orange × Teal） ────────────────────
N9 = RGBColor(0x0A, 0x23, 0x4E)  # 深海軍藍
N7 = RGBColor(0x1A, 0x50, 0x96)  # 中藍
N3 = RGBColor(0xBF, 0xD7, 0xFF)  # 淡藍
O7 = RGBColor(0xE5, 0x60, 0x00)  # 橘紅
O3 = RGBColor(0xFF, 0xC7, 0x8E)  # 淡橘
T7 = RGBColor(0x00, 0x97, 0xA7)  # 青色
T3 = RGBColor(0xB2, 0xEB, 0xF2)  # 淡青
WH = RGBColor(0xFF, 0xFF, 0xFF)  # 白
K1 = RGBColor(0x1A, 0x1A, 0x1A)  # 近黑
K5 = RGBColor(0x55, 0x55, 0x55)  # 中灰
K8 = RGBColor(0xF0, 0xF4, 0xFF)  # 極淡藍白（背景）
GD = RGBColor(0xEC, 0xEF, 0xF8)  # 淡灰藍（表格偶數列）

# ─── 尺寸常數 ──────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)
STRIPE = Inches(0.45)   # 左側色條寬度

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════
# 基礎繪圖工具
# ══════════════════════════════════════════════════════════

def bg(slide, color=WH):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, l, t, w, h, fill=N7, line=False, shape=1):
    """矩形 shape=1；橢圓 shape=9；圓角矩形 shape=5"""
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if not line:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = fill
        sp.line.width = Pt(1)
    return sp

def oval(slide, l, t, w, h, fill=N7, alpha_like=False):
    """橢圓（裝飾用）"""
    sp = slide.shapes.add_shape(9, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    return sp

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=K1,
        align=PP_ALIGN.LEFT, font=FT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def multi_para(slide, lines, l, t, w, h,
               size=18, color=K1, font=FT,
               spacing_pt=4, bullet_color=None, bullet_char="●"):
    """多段落文字方塊，支援 ● 色彩不同於文字"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if spacing_pt:
            p.space_before = Pt(spacing_pt)
        if bullet_color and line.strip():
            # 彩色項目符號
            br = p.add_run()
            br.text = bullet_char + "  "
            br.font.name = FT2
            br.font.size = Pt(size - 1)
            br.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb

def pnum(slide, n, total=20):
    """右下頁碼"""
    txt(slide, f"{n}  /  {total}",
        W - Inches(1.3), H - Inches(0.38),
        Inches(1.2), Inches(0.35),
        size=13, color=K5, align=PP_ALIGN.RIGHT, font=FT2)

def pdf_img(filename, page=0, dpi=180):
    doc = fitz.open(f"{DL}\\{filename}")
    pg  = doc[page]
    mat = fitz.Matrix(dpi/72, dpi/72)
    pix = pg.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
    data = pix.tobytes("png")
    doc.close()
    return data

def add_img(slide, data, l, t, w, h=None):
    import io as _io
    s = _io.BytesIO(data)
    return slide.shapes.add_picture(s, l, t, w, h)

def img_shadow(slide, l, t, w, h):
    """圖片陰影效果（深色偏移矩形）"""
    box(slide, l + Inches(0.08), t + Inches(0.08), w, h,
        fill=RGBColor(0xCC,0xCC,0xCC))

# ─── 左側色條 + 標題列（內容投影片通用）─────────────────────
def left_stripe_slide(title, sub=None, stripe_color=N9, bar_color=O7):
    """返回已設定左條＋標題的 slide"""
    s = prs.slides.add_slide(BLANK)
    bg(s, K8)

    # 左側深色條
    box(s, 0, 0, STRIPE, H, fill=stripe_color)
    # 左條上方橘色小方塊裝飾
    box(s, 0, 0, STRIPE, Inches(0.55), fill=bar_color)

    # 白色主內容區
    box(s, STRIPE, 0, W - STRIPE, H, fill=WH)

    # 標題底色帶
    box(s, STRIPE, 0, W - STRIPE, Inches(1.1), fill=K8)
    # 橘色細線
    box(s, STRIPE, Inches(1.1), W - STRIPE, Pt(3), fill=bar_color)

    # 標題文字
    txt(s, title,
        STRIPE + Inches(0.3), Inches(0.12),
        W - STRIPE - Inches(0.5), Inches(0.82),
        size=28, bold=True, color=N9, align=PP_ALIGN.LEFT)

    if sub:
        txt(s, sub,
            STRIPE + Inches(0.3), Inches(0.85),
            W - STRIPE - Inches(0.5), Inches(0.28),
            size=13, color=T7, align=PP_ALIGN.LEFT, italic=True, font=FT2)

    return s

# ─── 大標題數字（章節感）──────────────────────────────────
def big_num(slide, n, l, t, color=GD):
    txt(slide, str(n), l, t, Inches(1.2), Inches(1.4),
        size=96, bold=True, color=color, font=FT2)


# ══════════════════════════════════════════════════════════
# S01  封面
# ══════════════════════════════════════════════════════════
def s01():
    s = prs.slides.add_slide(BLANK)
    bg(s, N9)

    # 右上大圓裝飾
    oval(s, W - Inches(4.5), -Inches(3.0), Inches(7.0), Inches(7.0),
         fill=N7)
    oval(s, W - Inches(3.2), -Inches(1.8), Inches(4.8), Inches(4.8),
         fill=RGBColor(0x0D, 0x2D, 0x65))

    # 左下小橘圓
    oval(s, -Inches(1.0), H - Inches(2.5), Inches(4.0), Inches(4.0),
         fill=O7)
    oval(s, -Inches(0.2), H - Inches(1.8), Inches(2.5), Inches(2.5),
         fill=RGBColor(0xFF, 0x80, 0x00))

    # 橘色頂部細線
    box(s, 0, 0, W, Pt(6), fill=O7)

    # 學校名
    txt(s, "長庚大學  電機工程學系",
        Inches(1.0), Inches(0.8), Inches(8.0), Inches(0.55),
        size=18, color=N3, font=FT, align=PP_ALIGN.LEFT)

    # 主標題
    txt(s, "PWM IC  電路研製",
        Inches(1.0), Inches(1.45), Inches(9.5), Inches(1.0),
        size=46, bold=True, color=WH, font=FT, align=PP_ALIGN.LEFT)
    txt(s, "以 UC3845 BUCK 降壓轉換器為核心",
        Inches(1.0), Inches(1.15), Inches(9.5), Inches(0.55),
        size=20, color=O3, font=FT, align=PP_ALIGN.LEFT)

    # 英文副標題
    txt(s, "Design & Implementation of a BUCK Converter Based on UC3845 PWM IC",
        Inches(1.0), Inches(2.65), Inches(9.8), Inches(0.45),
        size=14, color=N3, font=FT2, align=PP_ALIGN.LEFT, italic=True)

    # 橘色分隔線
    box(s, Inches(1.0), Inches(3.2), Inches(2.5), Pt(3), fill=O7)

    # 教授 / 學生資訊卡
    box(s, Inches(1.0), Inches(3.45), Inches(4.4), Inches(1.55),
        fill=RGBColor(0x12, 0x33, 0x70))
    txt(s, "指導教授",
        Inches(1.15), Inches(3.55), Inches(4.0), Inches(0.4),
        size=13, color=O3, font=FT)
    txt(s, "曾聖有 教授",
        Inches(1.15), Inches(3.9), Inches(4.0), Inches(0.45),
        size=20, bold=True, color=WH, font=FT)

    box(s, Inches(5.7), Inches(3.45), Inches(5.2), Inches(1.55),
        fill=RGBColor(0x12, 0x33, 0x70))
    txt(s, "學生",
        Inches(5.85), Inches(3.55), Inches(4.8), Inches(0.4),
        size=13, color=O3, font=FT)
    txt(s, "周庭睿（B1221244）",
        Inches(5.85), Inches(3.9), Inches(4.8), Inches(0.42),
        size=19, bold=True, color=WH, font=FT)
    txt(s, "陳彥瑋（B1221213）",
        Inches(5.85), Inches(4.3), Inches(4.8), Inches(0.42),
        size=19, bold=True, color=WH, font=FT)

    # 日期
    txt(s, "中華民國 115 年 6 月 11 日",
        Inches(1.0), Inches(5.25), Inches(6.0), Inches(0.42),
        size=16, color=K5, font=FT)

    # 底部橘線
    box(s, 0, H - Pt(6), W, Pt(6), fill=O7)


# ══════════════════════════════════════════════════════════
# S02  目錄
# ══════════════════════════════════════════════════════════
def s02():
    s = left_stripe_slide("目  錄", "Table of Contents")
    bg_rect_l = STRIPE

    items = [
        ("01", "研究背景與目的"),
        ("02", "系統規格與研究流程"),
        ("03", "BUCK 轉換器工作原理"),
        ("04", "UC3845 PWM IC 介紹"),
        ("05", "峰值電流控制與補償網路"),
        ("06", "PSIM 電路模擬"),
        ("07", "Altium 電路圖 與 PCB 設計"),
        ("08", "量測與驗證結果"),
        ("09", "DSP 數位控制升級"),
        ("10", "結論與未來展望"),
    ]

    half = len(items) // 2
    cols = [(items[:half], Inches(0.7)), (items[half:], Inches(7.1))]

    for group, x_off in cols:
        for i, (num, title) in enumerate(group):
            y = Inches(1.25) + i * Inches(0.6)
            # 編號圓形
            oval(s, bg_rect_l + x_off, y + Inches(0.07),
                 Inches(0.42), Inches(0.42),
                 fill=N7 if int(num) % 2 == 1 else T7)
            txt(s, num,
                bg_rect_l + x_off, y + Inches(0.05),
                Inches(0.42), Inches(0.42),
                size=13, bold=True, color=WH,
                align=PP_ALIGN.CENTER, font=FT2)
            # 標題
            txt(s, title,
                bg_rect_l + x_off + Inches(0.55), y + Inches(0.06),
                Inches(5.6), Inches(0.44),
                size=19, color=K1, font=FT)
            # 底線
            if i < len(group) - 1:
                box(s, bg_rect_l + x_off, y + Inches(0.55),
                    Inches(6.1), Pt(1),
                    fill=RGBColor(0xDD, 0xDD, 0xEE))

    pnum(s, 2)


# ══════════════════════════════════════════════════════════
# S03  研究背景與目的
# ══════════════════════════════════════════════════════════
def s03():
    s = left_stripe_slide("研究背景與目的", "Background & Motivation")

    # 兩欄卡片
    for ci, (card_title, items, c_fill, t_fill) in enumerate([
        ("研究背景", [
            "電力電子廣泛應用於消費電子、電動車及再生能源",
            "DC-DC 轉換器是電源管理不可或缺的核心元件",
            "UC3845 專用 PWM IC 簡化控制電路設計複雜度",
            "DSP 數位控制取代類比電路，靈活性更高",
        ], N9, N3),
        ("研究目的", [
            "設計 12V～18V 輸入、5V / 2A 輸出 BUCK 轉換器",
            "以 UC3845 實現峰值電流控制（Peak Current Control）",
            "完成 PSIM 模擬→PCB 設計→製作→量測完整流程",
            "以 TMS320F28335 DSP 取代類比板，數位化升級",
        ], O7, O3),
    ]):
        x = STRIPE + Inches(0.3) + ci * Inches(6.3)
        # 卡片陰影
        box(s, x + Inches(0.08), Inches(1.25), Inches(5.9), Inches(5.8),
            fill=RGBColor(0xDD, 0xDD, 0xDD))
        # 卡片本體
        box(s, x, Inches(1.18), Inches(5.9), Inches(5.8), fill=WH)
        # 卡片頂色帶
        box(s, x, Inches(1.18), Inches(5.9), Inches(0.7), fill=c_fill)
        txt(s, card_title, x + Inches(0.2), Inches(1.22),
            Inches(5.5), Inches(0.62),
            size=22, bold=True, color=WH, font=FT)

        for j, item in enumerate(items):
            y = Inches(2.05) + j * Inches(1.1)
            oval(s, x + Inches(0.2), y + Inches(0.12),
                 Inches(0.28), Inches(0.28), fill=c_fill)
            txt(s, item,
                x + Inches(0.62), y,
                Inches(5.1), Inches(1.0),
                size=17, color=K1, font=FT, wrap=True)

    pnum(s, 3)


# ══════════════════════════════════════════════════════════
# S04  系統規格
# ══════════════════════════════════════════════════════════
def s04():
    s = left_stripe_slide("系統規格", "System Specifications")

    specs = [
        ("輸入電壓 Vin", "12V ～ 18V",  "直流輸入範圍"),
        ("輸出電壓 Vout","5V",           "目標輸出電壓"),
        ("輸出電流 Iout","2A",           "最大負載電流"),
        ("輸出功率",     "10W",          "Vout × Iout"),
        ("切換頻率 fs",  "100 kHz",      "Rt=7.5kΩ，Ct=2.2nF"),
        ("電感 L",       "1 mH",         "主電路儲能電感"),
        ("電容 C",       "4.7 mF",       "輸出濾波電容"),
        ("電流感測 Rs",  "0.1 Ω",        "峰值電流偵測"),
        ("控制模式",     "峰值電流控制", "Peak Current Control"),
    ]

    # 表頭
    hdr_x = [STRIPE+Inches(0.3), STRIPE+Inches(5.0), STRIPE+Inches(8.8)]
    hdr_w = [Inches(4.5), Inches(3.6), Inches(4.0)]
    hdr_labels = ["規格項目", "設計值", "說明"]
    for hx, hw, hl in zip(hdr_x, hdr_w, hdr_labels):
        box(s, hx, Inches(1.25), hw - Pt(4), Inches(0.5), fill=N9)
        txt(s, hl, hx + Inches(0.1), Inches(1.28),
            hw - Inches(0.15), Inches(0.46),
            size=17, bold=True, color=WH, font=FT,
            align=PP_ALIGN.CENTER)

    for i, (label, val, note) in enumerate(specs):
        y = Inches(1.8) + i * Inches(0.57)
        row_fill = WH if i % 2 == 0 else GD
        for hx, hw in zip(hdr_x, hdr_w):
            box(s, hx, y, hw - Pt(4), Inches(0.54), fill=row_fill)
        vals = [label, val, note]
        for hx, hw, v in zip(hdr_x, hdr_w, vals):
            fc = N9 if v == val else K1
            b  = (v == val)
            txt(s, v, hx + Inches(0.1), y + Pt(5),
                hw - Inches(0.2), Inches(0.5),
                size=17, bold=b, color=fc, font=FT,
                align=PP_ALIGN.CENTER if v == val else PP_ALIGN.LEFT)

    pnum(s, 4)


# ══════════════════════════════════════════════════════════
# S05  研究流程
# ══════════════════════════════════════════════════════════
def s05():
    s = left_stripe_slide("研究流程", "Research Methodology – 5 Steps")

    steps = [
        ("01", "PSIM\n模擬",     "建立電路模型\n驗證設計參數",     N9),
        ("02", "Altium\n電路設計","Schematic\nPCB Layout",          N7),
        ("03", "PCB\n製作焊接",  "送廠製板\n元件焊接",             T7),
        ("04", "量測\n驗證",     "示波器量測\n與模擬比對",          O7),
        ("05", "DSP\n升級",      "TMS320F28335\n數位 PID 控制",    RGBColor(0x6A,0x1B,0x9A)),
    ]

    step_w = Inches(2.22)
    gap    = Inches(0.12)

    for i, (num, title, desc, col) in enumerate(steps):
        x = STRIPE + Inches(0.3) + i * (step_w + gap)
        y0 = Inches(1.2)

        # 背景卡
        box(s, x + Inches(0.05), y0 + Inches(0.05),
            step_w, Inches(5.8), fill=RGBColor(0xCC,0xCC,0xCC))
        box(s, x, y0, step_w, Inches(5.8), fill=WH)

        # 頂色帶
        box(s, x, y0, step_w, Inches(1.6), fill=col)
        # 大數字
        txt(s, num, x, y0, step_w, Inches(0.8),
            size=44, bold=True, color=WH,
            align=PP_ALIGN.CENTER, font=FT2)
        # 步驟標題
        txt(s, title, x, y0 + Inches(0.75), step_w, Inches(0.85),
            size=17, bold=True, color=WH,
            align=PP_ALIGN.CENTER, font=FT)

        # 說明文字
        txt(s, desc,
            x + Inches(0.12), y0 + Inches(1.7),
            step_w - Inches(0.24), Inches(3.8),
            size=16, color=K1, font=FT,
            align=PP_ALIGN.CENTER, wrap=True)

        # 箭頭
        if i < len(steps) - 1:
            ax = x + step_w + Inches(0.01)
            txt(s, "▶", ax, y0 + Inches(0.65), gap + Inches(0.05), Inches(0.5),
                size=18, color=col, align=PP_ALIGN.CENTER, font=FT2)

    pnum(s, 5)


# ══════════════════════════════════════════════════════════
# S06  BUCK 原理
# ══════════════════════════════════════════════════════════
def s06():
    s = left_stripe_slide("BUCK 降壓轉換器工作原理",
                          "BUCK Converter Operating Principle")
    cx = STRIPE + Inches(0.3)

    # 大裝飾數字
    big_num(s, "06", W - Inches(1.5), H - Inches(1.6))

    # 左欄
    box(s, cx, Inches(1.25), Inches(5.8), Inches(0.48), fill=N9)
    txt(s, "ON 狀態 / OFF 狀態分析",
        cx + Inches(0.15), Inches(1.28), Inches(5.5), Inches(0.44),
        size=19, bold=True, color=WH, font=FT)

    on_items = [
        "MOSFET 導通，電感儲能",
        "電流路徑：Vin → L → C → Load",
        "電感電流線性上升",
        "ΔiL↑ = (Vin-Vo)×D / (L×fs)",
    ]
    off_items = [
        "MOSFET 截止，蕭特基二極體續流",
        "電感釋能，維持輸出電壓穩定",
        "電感電流線性下降",
        "ΔiL↓ = Vo×(1-D) / (L×fs)",
    ]
    for j, (label, items, c) in enumerate([
        ("● ON  狀態（Switch ON）",  on_items,  N7),
        ("● OFF 狀態（Switch OFF）", off_items, O7),
    ]):
        y = Inches(1.85) + j * Inches(2.6)
        txt(s, label, cx, y, Inches(5.8), Inches(0.44),
            size=18, bold=True, color=c, font=FT)
        for k, item in enumerate(items):
            txt(s, "  " + item,
                cx + Inches(0.1), y + Inches(0.5) + k * Inches(0.48),
                Inches(5.6), Inches(0.46),
                size=16, color=K1, font=FT)

    # 右欄公式卡片
    rx = STRIPE + Inches(6.4)
    box(s, rx + Inches(0.06), Inches(1.22), Inches(6.5), Inches(6.1),
        fill=RGBColor(0xCC,0xCC,0xCC))
    box(s, rx, Inches(1.17), Inches(6.5), Inches(6.1), fill=K8)
    box(s, rx, Inches(1.17), Inches(6.5), Inches(0.48), fill=T7)
    txt(s, "核心公式",
        rx + Inches(0.15), Inches(1.2), Inches(6.2), Inches(0.44),
        size=19, bold=True, color=WH, font=FT)

    formulas = [
        ("輸出電壓", "Vout = D × Vin", 26, N9),
        ("責任週期", "D = Vout/Vin = 5/18 ≈ 0.278（Vin=18V）", 16, K1),
        ("臨界電感", "Lmin = (1-D)Vout / (2·fs·Iout)", 16, K1),
        ("",         "= (1-0.278)×5 / (2×100k×2) ≈ 9μH", 16, K5),
        ("設計選用", "L = 1 mH  （遠大於臨界值，確保CCM）", 16, T7),
        ("漣波電壓", "ΔVo = (1-D)Vout / (8·L·C·fs²)", 16, K1),
        ("漣波估算", "≈ 5×10⁻³ V  （< 0.1% Vout）", 16, K5),
    ]
    for k, (label, formula, sz, fc) in enumerate(formulas):
        y = Inches(1.78) + k * Inches(0.68)
        if label:
            txt(s, label + "：",
                rx + Inches(0.2), y, Inches(1.6), Inches(0.55),
                size=14, color=K5, font=FT)
        txt(s, formula,
            rx + Inches(0.2 if not label else 1.8), y,
            Inches(6.0 if not label else 4.5), Inches(0.55),
            size=sz, bold=(sz > 20), color=fc, font=FT2 if sz > 20 else FT)

    pnum(s, 6)


# ══════════════════════════════════════════════════════════
# S07  UC3845
# ══════════════════════════════════════════════════════════
def s07():
    s = left_stripe_slide("UC3845 PWM IC 介紹",
                          "UC3845 PWM Controller – Internal Blocks")
    big_num(s, "07", W - Inches(1.5), H - Inches(1.6))

    modules = [
        ("振盪器",    "Oscillator",      "Rt=7.5kΩ，Ct=2.2nF\nfs = 1.72/(Rt×Ct) ≈ 104kHz", N9),
        ("誤差放大器","Error Amp",       "比較 Vfb(2.5V) 與回授電壓\n輸出補償訊號驅動 PWM 比較器", N7),
        ("PWM 比較器","PWM Comparator",  "電流感測 vs 誤差訊號\n決定每週期截止時間點", T7),
        ("驅動輸出",  "GATE Driver",     "推挽式輸出，搭配 IR2111\n驅動 IRLZ44N MOSFET", O7),
    ]

    for i, (zh, en, desc, c) in enumerate(modules):
        row = i // 2
        col = i %  2
        x = STRIPE + Inches(0.3) + col * Inches(6.3)
        y = Inches(1.2) + row * Inches(2.95)

        box(s, x + Inches(0.06), y + Inches(0.06),
            Inches(6.0), Inches(2.65), fill=RGBColor(0xCC,0xCC,0xCC))
        box(s, x, y, Inches(6.0), Inches(2.65), fill=WH)
        box(s, x, y, Inches(6.0), Inches(0.58), fill=c)
        box(s, x, y, Inches(0.08), Inches(2.65), fill=c)

        txt(s, zh, x + Inches(0.2), y + Inches(0.1),
            Inches(3.5), Inches(0.46),
            size=21, bold=True, color=WH, font=FT)
        txt(s, en, x + Inches(3.7), y + Inches(0.14),
            Inches(2.1), Inches(0.38),
            size=13, color=WH, font=FT2, italic=True)

        txt(s, desc, x + Inches(0.2), y + Inches(0.72),
            Inches(5.6), Inches(1.85),
            size=17, color=K1, font=FT, wrap=True)

    pnum(s, 7)


# ══════════════════════════════════════════════════════════
# S08  峰值電流控制 & 補償
# ══════════════════════════════════════════════════════════
def s08():
    s = left_stripe_slide("峰值電流控制與補償網路設計",
                          "Peak Current Control & Type II Compensation")
    big_num(s, "08", W - Inches(1.5), H - Inches(1.6))

    cx = STRIPE + Inches(0.3)

    # 左欄
    box(s, cx, Inches(1.25), Inches(5.9), Inches(0.48), fill=N9)
    txt(s, "峰值電流控制原理", cx + Inches(0.15), Inches(1.28),
        Inches(5.6), Inches(0.44), size=19, bold=True, color=WH, font=FT)
    pcc = [
        "電壓外迴路：誤差放大器設定峰值電流命令",
        "電流內迴路：感測電阻直接偵測電感峰值電流",
        "每週期重設：電流超過命令值 → GATE 截止",
        "優點：自動過電流保護、動態響應快",
        "D > 0.5 時需加入斜坡補償（Slope Compensation）",
    ]
    for j, p in enumerate(pcc):
        oval(s, cx + Inches(0.1), Inches(1.88) + j * Inches(0.55),
             Inches(0.22), Inches(0.22), fill=N7)
        txt(s, p, cx + Inches(0.45), Inches(1.85) + j * Inches(0.55),
            Inches(5.4), Inches(0.5), size=17, color=K1, font=FT)

    # 右欄卡片
    rx = STRIPE + Inches(6.5)
    box(s, rx + Inches(0.06), Inches(1.22),
        Inches(6.3), Inches(6.1), fill=RGBColor(0xCC,0xCC,0xCC))
    box(s, rx, Inches(1.17), Inches(6.3), Inches(6.1), fill=K8)
    box(s, rx, Inches(1.17), Inches(6.3), Inches(0.48), fill=O7)
    txt(s, "Type II 補償網路",
        rx + Inches(0.15), Inches(1.2), Inches(6.0), Inches(0.44),
        size=19, bold=True, color=WH, font=FT)

    comp = [
        ("補償目標",   "Phase Margin > 45°",                          T7),
        ("傳遞函數",   "Gc(s) = −Kc·(1+s/ωz) / [s·(1+s/ωp)]",       K1),
        ("零點 fz",    "fz = 1/(2π·Rc·Cp)\n  = 1/(2π×150k×4.7n) ≈ 226 Hz", N9),
        ("設計元件",   "Rc = 150kΩ，Cp = 4.7nF",                     N7),
        ("極點 fp",    "fp = 1/(2π·Rcomp·Ccomp)\n  = 1/(2π×150k×4.7n) ≈ 226Hz", O7),
        ("Dmax",       "≈ 0.5（UC3845 內部鎖存限制）",                K5),
    ]
    for k, (lbl, val, fc) in enumerate(comp):
        y = Inches(1.78) + k * Inches(0.73)
        box(s, rx + Inches(0.15), y, Inches(1.55), Inches(0.28), fill=fc)
        txt(s, lbl, rx + Inches(0.15), y, Inches(1.55), Inches(0.28),
            size=12, bold=True, color=WH, font=FT, align=PP_ALIGN.CENTER)
        txt(s, val, rx + Inches(1.85), y - Inches(0.04),
            Inches(4.3), Inches(0.7),
            size=15, color=K1, font=FT, wrap=True)

    pnum(s, 8)


# ══════════════════════════════════════════════════════════
# S09  PSIM 電路圖（半版圖片）
# ══════════════════════════════════════════════════════════
def s09():
    s = left_stripe_slide("PSIM 電路模擬", "PSIM Circuit Simulation Model")
    big_num(s, "09", W - Inches(1.5), H - Inches(1.6))

    cx = STRIPE + Inches(0.3)
    box(s, cx, Inches(1.25), Inches(4.0), Inches(0.48), fill=N9)
    txt(s, "模擬電路架構", cx + Inches(0.15), Inches(1.28),
        Inches(3.8), Inches(0.44), size=19, bold=True, color=WH, font=FT)

    sim_items = [
        "軟體：PSIM（電力電子模擬專用）",
        "輸入：Vi = 18V（最大值）",
        "MOSFET（理想開關）",
        "蕭特基二極體 1N5822",
        "電感 L = 1 mH",
        "電容 C = 4.7 mF",
        "UC3845 功能方塊模型",
        "Type II 補償網路",
        "回授分壓電阻設計 Vfb=2.5V",
    ]
    for j, item in enumerate(sim_items):
        oval(s, cx + Inches(0.1),
             Inches(1.88) + j * Inches(0.52),
             Inches(0.22), Inches(0.22), fill=T7)
        txt(s, item,
            cx + Inches(0.45), Inches(1.85) + j * Inches(0.52),
            Inches(3.5), Inches(0.48), size=16, color=K1, font=FT)

    # 圖片（右 3/4）
    img = pdf_img("PSIM SIMULATION.pdf", dpi=200)
    img_shadow(s, Inches(4.9), Inches(1.15), Inches(8.0), Inches(6.0))
    add_img(s, img, Inches(4.9), Inches(1.15), Inches(8.0))

    pnum(s, 9)


# ══════════════════════════════════════════════════════════
# S10  PSIM 結果
# ══════════════════════════════════════════════════════════
def s10():
    s = left_stripe_slide("PSIM 模擬結果", "Simulation Results – Steady State")
    big_num(s, "10", W - Inches(1.5), H - Inches(1.6))

    results = [
        ("輸出電壓 Vout", "4.976 V", "目標 5V，誤差 0.49%", N9),
        ("輸出電流 Iout", "1.991 A", "目標 2A，誤差 0.47%",  N7),
        ("回授電壓 Vfb",  "2.499 V", "參考 2.5V，誤差 0.04%", T7),
        ("整體誤差",      "< 0.5%",  "遠優於規格要求 5%",    O7),
    ]

    for i, (label, val, note, c) in enumerate(results):
        col = i % 2
        row = i // 2
        x = STRIPE + Inches(0.3) + col * Inches(6.3)
        y = Inches(1.2) + row * Inches(2.8)

        box(s, x + Inches(0.06), y + Inches(0.06),
            Inches(6.0), Inches(2.5), fill=RGBColor(0xCC,0xCC,0xCC))
        box(s, x, y, Inches(6.0), Inches(2.5), fill=WH)
        box(s, x, y, Inches(0.1), Inches(2.5), fill=c)

        txt(s, label, x + Inches(0.25), y + Inches(0.18),
            Inches(5.6), Inches(0.45), size=17, color=K5, font=FT)
        txt(s, val,   x + Inches(0.25), y + Inches(0.6),
            Inches(5.6), Inches(0.9),  size=42, bold=True, color=c, font=FT2)
        txt(s, note,  x + Inches(0.25), y + Inches(1.55),
            Inches(5.6), Inches(0.5),  size=15, color=K5, font=FT)

    # 結論 Banner
    box(s, STRIPE + Inches(0.3), Inches(6.95),
        W - STRIPE - Inches(0.6), Inches(0.48), fill=N9)
    txt(s, "✔  模擬驗證通過：三項輸出量誤差均 < 0.5%，設計參數正確，可進行 PCB 製作。",
        STRIPE + Inches(0.45), Inches(6.98),
        W - STRIPE - Inches(0.9), Inches(0.44),
        size=17, bold=True, color=WH, font=FT)

    pnum(s, 10)


# ══════════════════════════════════════════════════════════
# 通用：左文右圖投影片
# ══════════════════════════════════════════════════════════
def schematic_slide(n, title, sub, items, pdf_file, bullet_color=N7):
    s = left_stripe_slide(title, sub)
    big_num(s, f"{n:02d}", W - Inches(1.5), H - Inches(1.6))

    cx = STRIPE + Inches(0.3)
    box(s, cx, Inches(1.25), Inches(4.0), Inches(0.48), fill=N9)
    txt(s, "電路特點", cx + Inches(0.15), Inches(1.28),
        Inches(3.8), Inches(0.44), size=19, bold=True, color=WH, font=FT)

    for j, item in enumerate(items):
        oval(s, cx + Inches(0.1),
             Inches(1.88) + j * Inches(0.52),
             Inches(0.22), Inches(0.22), fill=bullet_color)
        txt(s, item,
            cx + Inches(0.45), Inches(1.85) + j * Inches(0.52),
            Inches(3.5), Inches(0.48), size=16, color=K1, font=FT)

    img = pdf_img(pdf_file, dpi=220)
    img_shadow(s, Inches(4.9), Inches(1.15), Inches(8.0), Inches(6.0))
    add_img(s, img, Inches(4.9), Inches(1.15), Inches(8.0))
    pnum(s, n)


# ══════════════════════════════════════════════════════════
# S11-S14  電路圖 & PCB
# ══════════════════════════════════════════════════════════
def s11():
    schematic_slide(11,
        "控制板電路圖（Altium Designer）",
        "Control Board Schematic – UC3845 & IR2111",
        [
            "核心 IC：UC3845（峰值電流控制）",
            "閘極驅動：IR2111（高側/低側）",
            "振盪器：Rt=7.5kΩ，Ct=2.2nF",
            "補償：Rcomp=150kΩ，Ccomp=4.7nF",
            "電流感測：Rsense=1kΩ（圖示值）",
            "回授分壓：Vfb=2.5V 設計點",
            "輔助電源：12V→UC3845 VCC",
        ], "控板.pdf", N9)

def s12():
    schematic_slide(12,
        "電力板電路圖（Altium Designer）",
        "Power Board Schematic – BUCK Main Circuit",
        [
            "功率開關：IRLZ44N N-MOSFET",
            "續流元件：1N5822 蕭特基二極體",
            "儲能電感：L = 1 mH",
            "濾波電容：C = 4.7 mF",
            "輸入：12V ～ 18V 直流",
            "輸出：5V / 2A（額定 10W）",
            "大電流走線加寬至 ≥ 2mm",
        ], "電板.pdf", T7)

def s13():
    schematic_slide(13,
        "PCB Layout — 控制板",
        "Control Board PCB Layout",
        [
            "雙層板設計（Top / Bottom）",
            "信號走線間距 ≥ 0.3mm",
            "UC3845 與補償元件集中佈局",
            "縮小感測迴路面積減少雜訊",
            "電源地與信號地分離",
            "THT 設計，便於手工焊接",
        ], "控板pcb.pdf", O7)

def s14():
    schematic_slide(14,
        "PCB Layout — 電力板",
        "Power Board PCB Layout",
        [
            "大電流走線加寬（≥ 2mm）",
            "MOSFET 散熱墊設計",
            "蕭特基二極體緊鄰 MOSFET",
            "電感電容靠近輸出端佈置",
            "鋪銅（Polygon Pour）強化接地",
            "測試點設置便於示波器量測",
        ], "電板pcb.pdf", RGBColor(0x6A,0x1B,0x9A))


# ══════════════════════════════════════════════════════════
# S15  PCB 製作成果（海報截圖）
# ══════════════════════════════════════════════════════════
def s15():
    s = left_stripe_slide("PCB 製作與焊接成果", "PCB Fabrication & Component Assembly")
    big_num(s, "15", W - Inches(1.5), H - Inches(1.6))

    cx = STRIPE + Inches(0.3)
    steps = [
        ("①", "Altium 匯出 Gerber 檔案",   N9),
        ("②", "送廠打樣（嘉立創）",          N7),
        ("③", "採購元件（立即電子）",         T7),
        ("④", "SMD / THT 焊接作業",          O7),
        ("⑤", "導通 & 短路測試",              N9),
        ("⑥", "低壓上電初步驗證",            N7),
    ]
    box(s, cx, Inches(1.25), Inches(4.0), Inches(0.48), fill=N9)
    txt(s, "製作流程", cx + Inches(0.15), Inches(1.28),
        Inches(3.8), Inches(0.44), size=19, bold=True, color=WH, font=FT)

    for j, (num, step, c) in enumerate(steps):
        y = Inches(1.88) + j * Inches(0.55)
        oval(s, cx + Inches(0.05), y + Inches(0.07),
             Inches(0.38), Inches(0.38), fill=c)
        txt(s, num, cx + Inches(0.05), y + Inches(0.04),
            Inches(0.38), Inches(0.38),
            size=15, bold=True, color=WH, font=FT, align=PP_ALIGN.CENTER)
        txt(s, step, cx + Inches(0.55), y + Inches(0.06),
            Inches(3.3), Inches(0.44), size=17, color=K1, font=FT)

    img = pdf_img("期末海報.pdf", dpi=150)
    img_shadow(s, Inches(4.9), Inches(1.15), Inches(8.0), Inches(6.0))
    add_img(s, img, Inches(4.9), Inches(1.15), Inches(8.0))
    pnum(s, 15)


# ══════════════════════════════════════════════════════════
# S16 & S17  量測波形
# ══════════════════════════════════════════════════════════
def waveform_slide(n, title, sub, items, pdf_file, bullet_c=N7):
    s = left_stripe_slide(title, sub)
    big_num(s, f"{n:02d}", W - Inches(1.5), H - Inches(1.6))

    cx = STRIPE + Inches(0.3)
    box(s, cx, Inches(1.25), Inches(4.0), Inches(0.48), fill=N9)
    txt(s, "量測說明", cx + Inches(0.15), Inches(1.28),
        Inches(3.8), Inches(0.44), size=19, bold=True, color=WH, font=FT)

    for j, item in enumerate(items):
        oval(s, cx + Inches(0.1),
             Inches(1.88) + j * Inches(0.55),
             Inches(0.22), Inches(0.22), fill=bullet_c)
        txt(s, item,
            cx + Inches(0.45), Inches(1.85) + j * Inches(0.55),
            Inches(3.5), Inches(0.5), size=16, color=K1, font=FT)

    img = pdf_img(pdf_file, dpi=220)
    img_shadow(s, Inches(4.9), Inches(1.15), Inches(8.0), Inches(6.0))
    add_img(s, img, Inches(4.9), Inches(1.15), Inches(8.0))
    pnum(s, n)

def s16():
    waveform_slide(16,
        "量測結果 — 輸出電壓波形",
        "Measurement – Output Voltage Waveform (Vout)",
        [
            "儀器：數位示波器  探棒 ×10",
            "量測點：輸出電容兩端 Vout",
            "穩態直流值 ≈ 5.2V",
            "目標 5V，誤差 4%",
            "疊加 100kHz PWM 漣波",
            "漣波 ΔVp-p ≈ 50mV",
            "穩態響應穩定，無異常振盪",
            "PSIM 模擬值：4.976V",
        ], "Vo.pdf", T7)

def s17():
    waveform_slide(17,
        "量測結果 — 輸出電流波形",
        "Measurement – Output Current Waveform (Iout)",
        [
            "量測方式：電流探棒（換算）",
            "量測點：輸出電感電流 iL",
            "穩態平均值 Iout ≈ 1.99A",
            "目標 2A，誤差 0.5%",
            "連續導通模式（CCM）",
            "電流漣波 ΔiL ≈ ±0.2A",
            "PSIM 模擬值：1.991A",
        ], "Io.pdf", O7)


# ══════════════════════════════════════════════════════════
# S18  DSP 升級
# ══════════════════════════════════════════════════════════
def s18():
    s = left_stripe_slide("DSP 數位控制升級 — TMS320F28335",
                          "Digital Control Upgrade – TMS320F28335 DSP Controller")
    big_num(s, "18", W - Inches(1.5), H - Inches(1.6))

    cx = STRIPE + Inches(0.3)

    # 左欄：三個小卡片
    cards = [
        ("升級動機", [
            "類比電路調參需更換元件",
            "數位 PID 可程式動態調整",
            "DSP 支援進階控制策略",
        ], N9),
        ("關鍵模組", [
            "ePWM：100kHz 切換訊號",
            "12-bit ADC：Vout 回授採樣",
            "150MHz CPU：即時 PID 運算",
        ], T7),
    ]
    for ci, (title, items, c) in enumerate(cards):
        y = Inches(1.2) + ci * Inches(2.9)
        box(s, cx + Inches(0.05), y + Inches(0.05),
            Inches(5.6), Inches(2.65), fill=RGBColor(0xCC,0xCC,0xCC))
        box(s, cx, y, Inches(5.6), Inches(2.65), fill=WH)
        box(s, cx, y, Inches(5.6), Inches(0.5), fill=c)
        box(s, cx, y, Inches(0.08), Inches(2.65), fill=c)
        txt(s, title, cx + Inches(0.2), y + Inches(0.08),
            Inches(5.2), Inches(0.44), size=20, bold=True, color=WH, font=FT)
        for j, item in enumerate(items):
            oval(s, cx + Inches(0.2), y + Inches(0.66) + j * Inches(0.65),
                 Inches(0.2), Inches(0.2), fill=c)
            txt(s, item, cx + Inches(0.52), y + Inches(0.63) + j * Inches(0.65),
                Inches(5.0), Inches(0.58), size=17, color=K1, font=FT)

    # 右欄：PID 公式卡片
    rx = STRIPE + Inches(6.2)
    box(s, rx + Inches(0.06), Inches(1.22),
        Inches(6.7), Inches(6.2), fill=RGBColor(0xCC,0xCC,0xCC))
    box(s, rx, Inches(1.17), Inches(6.7), Inches(6.2), fill=K8)
    box(s, rx, Inches(1.17), Inches(6.7), Inches(0.5), fill=O7)
    txt(s, "數位 PID 控制架構",
        rx + Inches(0.2), Inches(1.2), Inches(6.4), Inches(0.46),
        size=20, bold=True, color=WH, font=FT)

    pid_blocks = [
        ("離散 PID 方程式（後退差分法）", None, 17, N9, True),
        ("u[k] = u[k-1]", None, 22, N9, True),
        ("  + Kp·(e[k] − e[k-1])", None, 18, K1, False),
        ("  + Ki·Ts·e[k]", None, 18, K1, False),
        ("  + Kd/Ts·(e[k] − 2e[k-1] + e[k-2])", None, 18, K1, False),
        ("", None, 8, WH, False),
        ("e[k] = Vref − Vout_sampled", "誤差訊號", 16, K1, False),
        ("u[k] = PWM 責任週期命令", "控制輸出", 16, K1, False),
        ("Ts = 10μs（fs=100kHz）", "採樣週期", 16, K1, False),
        ("TBPRD = 1500（150MHz/100kHz）", "ePWM 設定", 16, T7, False),
    ]
    tb = s.shapes.add_textbox(rx + Inches(0.2), Inches(1.78),
                               Inches(6.3), Inches(5.4))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (line, label, sz, fc, bld) in enumerate(pid_blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if label:
            lr = p.add_run()
            lr.text = f"[{label}]  "
            lr.font.name = FT2
            lr.font.size = Pt(11)
            lr.font.color.rgb = K5
        r = p.add_run()
        r.text = line
        r.font.name = FT2
        r.font.size = Pt(sz)
        r.font.bold = bld
        r.font.color.rgb = fc

    pnum(s, 18)


# ══════════════════════════════════════════════════════════
# S19  結果比較總表
# ══════════════════════════════════════════════════════════
def s19():
    s = left_stripe_slide("量測結果總整理與比較",
                          "Summary: Simulation vs. Measurement Results")
    big_num(s, "19", W - Inches(1.5), H - Inches(1.6))

    table = [
        ("量測項目",      "設計規格",  "PSIM 模擬",  "實際量測",  "誤差評估"),
        ("輸出電壓 Vout", "5.0 V",     "4.976 V",   "5.2 V",    "4.0%   ✔"),
        ("輸出電流 Iout", "2.0 A",     "1.991 A",   "1.99 A",   "0.5%   ✔"),
        ("切換頻率 fs",   "100 kHz",   "100 kHz",   "~100 kHz", "< 1%  ✔"),
        ("電壓漣波 ΔV",   "< 100mV",   "—",         "~50 mV",   "符合   ✔"),
        ("控制模式",      "峰值電流",  "峰值電流",  "峰值電流", "—     ✔"),
        ("電路效率 η",    "> 85%",     "—",         "~87%",     "符合   ✔"),
    ]

    col_x = [STRIPE+Inches(0.3), STRIPE+Inches(3.1),
             STRIPE+Inches(5.6), STRIPE+Inches(7.9), STRIPE+Inches(10.2)]
    col_w = [Inches(2.65), Inches(2.35), Inches(2.15), Inches(2.15), Inches(2.35)]
    row_h = Inches(0.62)

    for r, row in enumerate(table):
        y = Inches(1.22) + r * row_h
        for c, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
            is_hdr = (r == 0)
            if is_hdr:
                fill = N9; fc = WH
            elif r % 2 == 1:
                fill = GD; fc = K1
            else:
                fill = WH; fc = K1

            box(s, cx, y, cw - Pt(3), row_h - Pt(2), fill=fill)

            # 最後欄特殊顏色
            if not is_hdr and c == len(row) - 1:
                box(s, cx, y, cw - Pt(3), row_h - Pt(2),
                    fill=RGBColor(0xE8, 0xF5, 0xE9))
                fc = RGBColor(0x1B, 0x5E, 0x20)

            bold = is_hdr or (c == 0 and not is_hdr)
            alg  = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            txt(s, cell, cx + Inches(0.1), y + Pt(5),
                cw - Inches(0.15), row_h,
                size=16 if not is_hdr else 17,
                bold=bold, color=fc, font=FT, align=alg)

    # 結論
    box(s, STRIPE + Inches(0.3), Inches(5.62),
        W - STRIPE - Inches(0.6), Inches(0.52), fill=N9)
    txt(s, "✔  所有量測結果符合設計規格，模擬→設計→製作→量測完整驗證流程已成功執行。",
        STRIPE + Inches(0.45), Inches(5.65),
        W - STRIPE - Inches(0.9), Inches(0.46),
        size=17, bold=True, color=WH, font=FT)

    pnum(s, 19)


# ══════════════════════════════════════════════════════════
# S20  結論與展望
# ══════════════════════════════════════════════════════════
def s20():
    s = prs.slides.add_slide(BLANK)
    bg(s, N9)

    # 頂部橘色線
    box(s, 0, 0, W, Pt(6), fill=O7)
    box(s, 0, H - Pt(6), W, Pt(6), fill=O7)

    # 右上裝飾圓
    oval(s, W - Inches(4.0), -Inches(2.2), Inches(5.5), Inches(5.5),
         fill=N7)
    oval(s, W - Inches(2.8), -Inches(1.2), Inches(3.5), Inches(3.5),
         fill=RGBColor(0x0D, 0x2D, 0x65))

    # 左下橘圓
    oval(s, -Inches(0.8), H - Inches(2.2), Inches(3.5), Inches(3.5),
         fill=O7)

    txt(s, "結論與未來展望",
        Inches(0.7), Inches(0.15), Inches(9.0), Inches(0.75),
        size=36, bold=True, color=WH, font=FT)

    # 兩欄卡片
    for ci, (card_title, items, c) in enumerate([
        ("研究結論", [
            "成功設計並製作以 UC3845 為核心的 BUCK 降壓轉換器",
            "規格達標：Vin 12～18V、Vout 5V/2A、fs 100kHz",
            "PSIM 模擬誤差 < 0.5%，設計參數驗證正確",
            "實測 Vout≈5.2V、Iout≈1.99A，誤差均在 5% 以內",
            "PCB 完整流程：Altium 設計→送廠→焊接→量測",
            "完成 TMS320F28335 數位 PID 升級架構設計",
        ], N7),
        ("未來展望", [
            "優化 MOSFET 驅動，降低切換損耗提升效率",
            "加入軟切換（ZVS）技術降低電磁干擾",
            "擴展至雙向 DC-DC（BUCK-BOOST）架構",
            "完善 DSP 數位 PID：前饋補償提升動態響應",
            "設計閉環 Bode Plot 量測，驗證相位裕度",
            "應用至再生能源儲能系統電池充電管理",
        ], T7),
    ]):
        x = Inches(0.5) + ci * Inches(6.4)
        box(s, x + Inches(0.07), Inches(1.05),
            Inches(6.1), Inches(5.85), fill=RGBColor(0xCC,0xCC,0xCC))
        box(s, x, Inches(0.98), Inches(6.1), Inches(5.85),
            fill=RGBColor(0x12, 0x33, 0x70))
        box(s, x, Inches(0.98), Inches(6.1), Inches(0.55), fill=c)
        txt(s, card_title, x + Inches(0.2), Inches(1.01),
            Inches(5.7), Inches(0.5),
            size=22, bold=True, color=WH, font=FT)

        for j, item in enumerate(items):
            oval(s, x + Inches(0.2), Inches(1.65) + j * Inches(0.82),
                 Inches(0.24), Inches(0.24), fill=c)
            txt(s, item,
                x + Inches(0.6), Inches(1.62) + j * Inches(0.82),
                Inches(5.3), Inches(0.78),
                size=16, color=N3, font=FT, wrap=True)

    # 致謝
    txt(s, "衷心感謝  曾聖有 教授  悉心指導",
        Inches(0.5), H - Inches(0.55),
        W - Inches(1.0), Inches(0.45),
        size=18, color=O3, font=FT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# 執行
# ══════════════════════════════════════════════════════════
print("正在生成投影片（v2 現代化設計）...")
s01(); print("  [ 1/20] done")
s02(); print("  [ 2/20] done")
s03(); print("  [ 3/20] done")
s04(); print("  [ 4/20] done")
s05(); print("  [ 5/20] done")
s06(); print("  [ 6/20] done")
s07(); print("  [ 7/20] done")
s08(); print("  [ 8/20] done")
s09(); print("  [ 9/20] done")
s10(); print("  [10/20] done")
s11(); print("  [11/20] done")
s12(); print("  [12/20] done")
s13(); print("  [13/20] done")
s14(); print("  [14/20] done")
s15(); print("  [15/20] done")
s16(); print("  [16/20] done")
s17(); print("  [17/20] done")
s18(); print("  [18/20] done")
s19(); print("  [19/20] done")
s20(); print("  [20/20] done")

prs.save(OUT_UPL)
print(f"\nSaved upload copy: {OUT_UPL}")
try:
    prs.save(OUT)
    print(f"Saved desktop: {OUT}")
except Exception as e:
    print(f"Desktop save skipped (file open?): {e}")
