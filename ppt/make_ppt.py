"""
期末專題 PPT 生成腳本
主題：以 UC3845 PWM IC 為核心之 BUCK 降壓轉換器設計
"""

import fitz  # PyMuPDF
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── 路徑設定 ─────────────────────────────────────────────
DOWNLOADS = r"C:\Users\User\Downloads"
OUTPUT    = r"C:\Users\User\Desktop\期末專題PPT.pptx"
FONT      = "標楷體"

# ── 色彩定義 ─────────────────────────────────────────────
C_DARK_BLUE  = RGBColor(0x1F, 0x38, 0x64)
C_MID_BLUE   = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)
C_ORANGE     = RGBColor(0xC5, 0x50, 0x0C)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY       = RGBColor(0x40, 0x40, 0x40)
C_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

# ── 簡報尺寸 16:9 ─────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # 完全空白版型


# ══════════════════════════════════════════════════════════
# 工具函式
# ══════════════════════════════════════════════════════════

def bg(slide, color=C_WHITE):
    """設定投影片背景色"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color=C_MID_BLUE, line=False):
    """加入色塊矩形"""
    sp = slide.shapes.add_shape(1, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_color
    if not line:
        sp.line.fill.background()
    return sp

def textbox(slide, text, l, t, w, h,
            size=18, bold=False, color=C_BLACK,
            align=PP_ALIGN.LEFT, font=FONT, wrap=True):
    """加入文字方塊"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def textbox_lines(slide, lines, l, t, w, h,
                  size=18, color=C_BLACK, font=FONT,
                  line_spacing=1.2, bold_first=False):
    """多行文字方塊，支援清單"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
    return tb

def add_image_from_bytes(slide, img_bytes, l, t, w, h=None):
    """將圖片位元組插入投影片"""
    stream = io.BytesIO(img_bytes)
    if h:
        pic = slide.shapes.add_picture(stream, l, t, w, h)
    else:
        pic = slide.shapes.add_picture(stream, l, t, w)
    return pic

def pdf_to_img(filename, page=0, dpi=180):
    """從 PDF 擷取指定頁面為 PNG bytes"""
    path = f"{DOWNLOADS}\\{filename}"
    doc  = fitz.open(path)
    pg   = doc[page]
    mat  = fitz.Matrix(dpi/72, dpi/72)
    pix  = pg.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
    data = pix.tobytes("png")
    doc.close()
    return data

def title_bar(slide, title_text, sub_text=None, bar_h=Inches(1.25)):
    """標準投影片頂部深藍標題列"""
    rect(slide, 0, 0, W, bar_h, C_DARK_BLUE)
    textbox(slide, title_text,
            Inches(0.4), Inches(0.1), Inches(12.5), bar_h - Inches(0.1),
            size=32, bold=True, color=C_WHITE,
            align=PP_ALIGN.LEFT)
    if sub_text:
        textbox(slide, sub_text,
                Inches(0.4), Inches(0.8), Inches(12.5), Inches(0.45),
                size=16, color=C_LIGHT_BLUE, align=PP_ALIGN.LEFT)
    return bar_h

def divider(slide, y):
    """水平分隔線"""
    ln = slide.shapes.add_shape(1, Inches(0.4), y, W - Inches(0.8), Pt(2))
    ln.fill.solid()
    ln.fill.fore_color.rgb = C_MID_BLUE
    ln.line.fill.background()

def bullet_box(slide, items, l, t, w, h, size=20, indent="   ● "):
    """項目符號清單"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = indent + item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = C_BLACK
    return tb

def page_num(slide, n):
    """右下角頁碼"""
    textbox(slide, str(n),
            W - Inches(0.6), H - Inches(0.4),
            Inches(0.5), Inches(0.35),
            size=14, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════
# 投影片 1：封面
# ══════════════════════════════════════════════════════════
def slide01():
    s = prs.slides.add_slide(BLANK)
    bg(s, C_DARK_BLUE)

    # 頂部橘色裝飾條
    rect(s, 0, 0, W, Inches(0.18), C_ORANGE)
    # 底部橘色裝飾條
    rect(s, 0, H - Inches(0.18), W, Inches(0.18), C_ORANGE)

    # 中央白色背景卡片
    rect(s, Inches(0.8), Inches(0.6), Inches(11.73), Inches(6.0), C_WHITE)

    # 學校名稱
    textbox(s, "長庚大學 電機工程學系",
            Inches(1.0), Inches(0.8), Inches(11.33), Inches(0.55),
            size=22, color=C_MID_BLUE, align=PP_ALIGN.CENTER)

    # 主標題
    textbox(s, "以 UC3845 PWM IC 為核心之",
            Inches(1.0), Inches(1.5), Inches(11.33), Inches(0.85),
            size=36, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)
    textbox(s, "BUCK 降壓轉換器電路研製",
            Inches(1.0), Inches(2.25), Inches(11.33), Inches(0.85),
            size=36, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

    # 英文副標題
    textbox(s, "Design and Implementation of a BUCK Converter Based on UC3845 PWM IC",
            Inches(1.0), Inches(3.15), Inches(11.33), Inches(0.5),
            size=15, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 分隔線
    divider(s, Inches(3.75))

    # 指導教授
    textbox(s, "指導教授：曾聖有 教授",
            Inches(1.0), Inches(3.95), Inches(5.3), Inches(0.45),
            size=20, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

    # 學生
    textbox(s, "學　　生：周庭睿（B1221244）\n　　　　　陳彥瑋（B1221213）",
            Inches(6.8), Inches(3.95), Inches(5.3), Inches(0.8),
            size=20, color=C_DARK_BLUE, align=PP_ALIGN.LEFT)

    # 日期
    textbox(s, "中華民國 115 年 6 月",
            Inches(1.0), Inches(5.1), Inches(11.33), Inches(0.45),
            size=18, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# 投影片 2：目錄
# ══════════════════════════════════════════════════════════
def slide02():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "目  錄", "Outline")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    contents = [
        ("一", "研究背景與目的"),
        ("二", "系統規格與研究流程"),
        ("三", "BUCK 轉換器工作原理"),
        ("四", "UC3845 PWM IC 介紹"),
        ("五", "峰值電流控制與補償網路設計"),
        ("六", "PSIM 電路模擬"),
        ("七", "電路圖與 PCB 設計（Altium Designer）"),
        ("八", "量測與驗證結果"),
        ("九", "DSP 數位控制升級（TMS320F28335）"),
        ("十", "結論與未來展望"),
    ]

    for i, (num, title) in enumerate(contents):
        y = Inches(1.45) + i * Inches(0.58)
        # 章節號碼藍色方塊
        rect(s, Inches(0.5), y, Inches(0.55), Inches(0.46), C_MID_BLUE)
        textbox(s, num, Inches(0.5), y, Inches(0.55), Inches(0.46),
                size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        textbox(s, title, Inches(1.2), y + Inches(0.03),
                Inches(11.0), Inches(0.46),
                size=20, color=C_BLACK)
        # 淡色分隔線
        if i < len(contents) - 1:
            ln = s.shapes.add_shape(1, Inches(1.2), y + Inches(0.5),
                                    Inches(11.0), Pt(1))
            ln.fill.solid(); ln.fill.fore_color.rgb = C_LIGHT_BLUE
            ln.line.fill.background()

    page_num(s, 2)


# ══════════════════════════════════════════════════════════
# 投影片 3：研究背景與目的
# ══════════════════════════════════════════════════════════
def slide03():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "研究背景與目的", "Background & Motivation")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    # 左欄：背景
    rect(s, Inches(0.4), Inches(1.45), Inches(5.9), Inches(0.5), C_MID_BLUE)
    textbox(s, "研究背景", Inches(0.4), Inches(1.45), Inches(5.9), Inches(0.5),
            size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    bg_items = [
        "電力電子技術廣泛應用於消費電子、電動車及再生能源等領域",
        "DC-DC 轉換器（如 BUCK）是電源管理不可或缺的核心元件",
        "專用 PWM IC（UC3845）可簡化控制電路設計，降低開發複雜度",
        "DSP 數位控制器取代類比控制，具更高靈活性與精準度",
    ]
    bullet_box(s, bg_items, Inches(0.4), Inches(2.0), Inches(5.9), Inches(3.5), size=18)

    # 右欄：目的
    rect(s, Inches(7.0), Inches(1.45), Inches(5.9), Inches(0.5), C_ORANGE)
    textbox(s, "研究目的", Inches(7.0), Inches(1.45), Inches(5.9), Inches(0.5),
            size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    obj_items = [
        "設計 12V～18V 輸入、5V/2A 輸出之 BUCK 降壓轉換器",
        "以 UC3845 PWM IC 實現峰值電流控制（Peak Current Control）",
        "完成 PSIM 模擬、PCB 設計製作與實體量測驗證",
        "以 TMS320F28335 DSP 取代類比控制板，實現數位化升級",
    ]
    bullet_box(s, obj_items, Inches(7.0), Inches(2.0), Inches(5.9), Inches(3.5), size=18)

    # 底部橘色裝飾條
    rect(s, 0, H - Inches(0.5), W, Inches(0.5), C_LIGHT_GRAY)
    textbox(s, "長庚大學 電機工程學系 ｜ 指導教授：曾聖有",
            Inches(0.4), H - Inches(0.45), Inches(10), Inches(0.4),
            size=13, color=C_GRAY)
    page_num(s, 3)


# ══════════════════════════════════════════════════════════
# 投影片 4：系統規格
# ══════════════════════════════════════════════════════════
def slide04():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "系統規格", "System Specifications")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    # 規格表格用形狀模擬
    specs = [
        ("參數", "設計值", "說明"),
        ("輸入電壓 Vin", "12V ～ 18V", "直流輸入範圍"),
        ("輸出電壓 Vout", "5V", "目標輸出電壓"),
        ("輸出電流 Iout", "2A", "最大負載電流"),
        ("輸出功率", "10W", "Vout × Iout"),
        ("切換頻率 fs", "100 kHz", "Rt=7.5kΩ，Ct=2.2nF"),
        ("電感值 L", "1 mH", "主電路儲能電感"),
        ("輸出電容 C", "4.7 mF", "輸出濾波電容"),
        ("電流感測電阻 Rsense", "1 kΩ → 0.1Ω", "峰值電流偵測"),
    ]

    col_x = [Inches(0.4), Inches(5.2), Inches(9.5)]
    col_w = [Inches(4.6), Inches(4.1), Inches(3.6)]
    row_h = Inches(0.52)

    for r, row in enumerate(specs):
        y = Inches(1.45) + r * row_h
        for c, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
            is_header = (r == 0)
            fill = C_DARK_BLUE if is_header else (C_LIGHT_BLUE if r % 2 == 0 else C_WHITE)
            fc   = C_WHITE if is_header else C_BLACK
            rect(s, x, y, w - Inches(0.06), row_h - Pt(2), fill)
            textbox(s, cell, x + Inches(0.1), y + Pt(4),
                    w - Inches(0.2), row_h,
                    size=18 if not is_header else 20,
                    bold=is_header, color=fc,
                    align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)

    page_num(s, 4)


# ══════════════════════════════════════════════════════════
# 投影片 5：研究流程
# ══════════════════════════════════════════════════════════
def slide05():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "研究流程", "Research Flowchart")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    steps = [
        ("STEP 1", "PSIM 電路模擬",
         "使用 PSIM 軟體建立 BUCK 轉換器電路模型\n驗證設計參數（Vo=4.98V, Io=1.99A）"),
        ("STEP 2", "Altium Designer 設計",
         "完成電路原理圖（Schematic）\n與 PCB Layout 佈線設計"),
        ("STEP 3", "PCB 製作與焊接",
         "送廠製板、採購元件\n完成 SMD/Through-hole 焊接"),
        ("STEP 4", "量測與驗證",
         "示波器量測輸出電壓/電流波形\n與模擬結果比對（誤差 < 5%）"),
        ("STEP 5", "DSP 數位控制升級",
         "以 TMS320F28335 ePWM + ADC\n實現數位 PID 補償控制"),
    ]

    for i, (step, title, desc) in enumerate(steps):
        x = Inches(0.4) + i * Inches(2.55)
        y_box = Inches(1.5)

        # 步驟藍色方塊
        rect(s, x, y_box, Inches(2.3), Inches(0.6), C_DARK_BLUE)
        textbox(s, step, x, y_box, Inches(2.3), Inches(0.6),
                size=16, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

        # 步驟標題
        rect(s, x, y_box + Inches(0.6), Inches(2.3), Inches(0.6), C_MID_BLUE)
        textbox(s, title, x, y_box + Inches(0.6), Inches(2.3), Inches(0.6),
                size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # 說明白色卡片
        rect(s, x, y_box + Inches(1.2), Inches(2.3), Inches(3.5), C_LIGHT_GRAY)
        textbox(s, desc, x + Inches(0.1), y_box + Inches(1.35),
                Inches(2.1), Inches(3.3),
                size=16, color=C_BLACK, align=PP_ALIGN.LEFT)

        # 箭頭（除最後一個）
        if i < len(steps) - 1:
            ax = x + Inches(2.3) + Inches(0.05)
            ay = y_box + Inches(0.95)
            textbox(s, "▶", ax, ay, Inches(0.25), Inches(0.5),
                    size=22, color=C_ORANGE, align=PP_ALIGN.CENTER)

    page_num(s, 5)


# ══════════════════════════════════════════════════════════
# 投影片 6：BUCK 轉換器工作原理
# ══════════════════════════════════════════════════════════
def slide06():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "BUCK 降壓轉換器工作原理", "BUCK Converter Operating Principle")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    # 左側說明
    textbox(s, "工作模式說明", Inches(0.4), Inches(1.45), Inches(6.0), Inches(0.45),
            size=22, bold=True, color=C_DARK_BLUE)
    divider(s, Inches(1.95))

    on_items = [
        "開關導通（Switch ON）",
        "MOSFET 導通，電感儲能",
        "電流路徑：Vin → L → C → Load",
        "電感電流線性上升",
        "iL(t) = iL(0) + (Vin-Vo)/L × t",
    ]
    textbox(s, "● ON 狀態", Inches(0.4), Inches(2.05), Inches(5.8), Inches(0.45),
            size=19, bold=True, color=C_MID_BLUE)
    bullet_box(s, on_items[1:], Inches(0.4), Inches(2.5), Inches(5.8), Inches(1.8),
               size=17, indent="   - ")

    off_items = [
        "開關截止（Switch OFF）",
        "二極體（蕭特基）續流",
        "電感釋能維持輸出電壓穩定",
        "電感電流線性下降",
        "iL(t) = iL(0) - Vo/L × t",
    ]
    textbox(s, "● OFF 狀態", Inches(0.4), Inches(4.0), Inches(5.8), Inches(0.45),
            size=19, bold=True, color=C_ORANGE)
    bullet_box(s, off_items[1:], Inches(0.4), Inches(4.45), Inches(5.8), Inches(1.8),
               size=17, indent="   - ")

    # 右側：輸出電壓與責任週期公式
    rect(s, Inches(6.8), Inches(1.45), Inches(6.1), Inches(5.5), C_LIGHT_GRAY)

    formulas = [
        ("輸出電壓與責任週期", 22, True, C_DARK_BLUE),
        ("", 10, False, C_WHITE),
        ("Vout = D × Vin", 30, True, C_MID_BLUE),
        ("", 10, False, C_WHITE),
        ("D = Vout / Vin = 5 / 12 ≈ 0.417", 19, False, C_BLACK),
        ("（Vin = 12V 時）", 16, False, C_GRAY),
        ("", 10, False, C_WHITE),
        ("臨界電感量（CCM/DCM 邊界）", 20, True, C_DARK_BLUE),
        ("", 8, False, C_WHITE),
        ("Lmin = (1-D)×Vout / (2×fs×Iout)", 20, False, C_BLACK),
        ("     = (1-0.417)×5 / (2×100k×2)", 18, False, C_GRAY),
        ("     ≈ 7.29 μH（設計使用 1 mH）", 18, False, C_BLACK),
        ("", 10, False, C_WHITE),
        ("輸出電壓漣波", 20, True, C_DARK_BLUE),
        ("", 8, False, C_WHITE),
        ("ΔVo = (1-D)×Vout / (8×L×C×fs²)", 18, False, C_BLACK),
    ]

    tb = slide.shapes.add_textbox if False else None
    txBox = s.shapes.add_textbox(Inches(7.0), Inches(1.55), Inches(5.7), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, sz, bld, clr) in enumerate(formulas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = txt
        run.font.name = FONT
        run.font.size = Pt(sz)
        run.font.bold = bld
        run.font.color.rgb = clr

    page_num(s, 6)


# ══════════════════════════════════════════════════════════
# 投影片 7：UC3845 PWM IC 介紹
# ══════════════════════════════════════════════════════════
def slide07():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "UC3845 PWM IC 介紹", "UC3845 PWM Controller")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    # 左側：IC 功能模組
    rect(s, Inches(0.4), Inches(1.45), Inches(6.0), Inches(0.45), C_DARK_BLUE)
    textbox(s, "IC 內部功能模組", Inches(0.4), Inches(1.45), Inches(6.0), Inches(0.45),
            size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    modules = [
        ("振盪器（Oscillator）",
         f"頻率由外部 Rt、Ct 決定\nfs = 1.72 / (Rt × Ct)\n設計：Rt=7.5kΩ，Ct=2.2nF → fs≈104kHz"),
        ("誤差放大器（Error Amp）",
         "比較參考電壓（2.5V）與回授電壓\n誤差訊號驅動 PWM 比較器"),
        ("PWM 比較器",
         "電流感測訊號與誤差訊號比較\n決定 GATE 截止時間點"),
        ("驅動輸出（GATE）",
         "推挽式輸出，驅動 MOSFET\n搭配 IR2111 閘極驅動 IC"),
    ]

    for i, (mod, desc) in enumerate(modules):
        y = Inches(2.0) + i * Inches(1.2)
        # 模組編號
        rect(s, Inches(0.4), y, Inches(0.4), Inches(0.4),
             C_MID_BLUE if i % 2 == 0 else C_ORANGE)
        textbox(s, str(i+1), Inches(0.4), y, Inches(0.4), Inches(0.4),
                size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        textbox(s, mod, Inches(0.95), y, Inches(5.4), Inches(0.4),
                size=18, bold=True, color=C_DARK_BLUE)
        textbox(s, desc, Inches(0.95), y + Inches(0.4), Inches(5.4), Inches(0.75),
                size=15, color=C_GRAY)

    # 右側：頻率與補償參數
    rect(s, Inches(6.8), Inches(1.45), Inches(6.1), Inches(5.5), C_LIGHT_GRAY)
    textbox(s, "關鍵設計參數", Inches(6.9), Inches(1.55), Inches(5.9), Inches(0.45),
            size=22, bold=True, color=C_DARK_BLUE)
    divider(s, Inches(2.05))

    params = [
        "振盪器設定",
        f"  Rt = 7.5 kΩ，Ct = 2.2 nF",
        f"  fs = 1.72 / (7.5k × 2.2n) ≈ 104 kHz",
        "",
        "補償網路（Type II）",
        "  Rc = 150 kΩ（補償零點電阻）",
        "  Cp = 4.7 nF（補償電容）",
        f"  零點 fz = 1/(2π×Rc×Cp) ≈ 226 Hz",
        "",
        "電流感測",
        "  Rsense = 0.1 Ω（實際量測用）",
        "  峰值電流限制 = 1A/div",
        "",
        "最大責任週期",
        "  Dmax ≈ 0.5（50%，內部鎖存）",
    ]
    bullet_box(s, params, Inches(6.9), Inches(2.1), Inches(5.9), Inches(4.8),
               size=17, indent="")

    page_num(s, 7)


# ══════════════════════════════════════════════════════════
# 投影片 8：峰值電流控制與補償網路
# ══════════════════════════════════════════════════════════
def slide08():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "峰值電流控制與補償網路設計",
                      "Peak Current Control & Compensation Network")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    # 左欄
    textbox(s, "峰值電流控制原理", Inches(0.4), Inches(1.45), Inches(6.2), Inches(0.45),
            size=22, bold=True, color=C_DARK_BLUE)
    divider(s, Inches(1.95))
    pcc_items = [
        "外迴路：電壓誤差放大器設定峰值電流命令",
        "內迴路：電流感測比較器直接偵測電感峰值電流",
        "每個切換週期重設：電流超過命令值即截止 GATE",
        "優點：自動過電流保護、動態響應快、次諧波振盪需補償",
        "次諧波振盪（D > 0.5）：加入 50% 斜坡補償（Slope Compensation）",
    ]
    bullet_box(s, pcc_items, Inches(0.4), Inches(2.0), Inches(6.2), Inches(3.5), size=17)

    # 右欄：補償網路
    rect(s, Inches(7.0), Inches(1.45), Inches(5.9), Inches(5.5), C_LIGHT_GRAY)
    textbox(s, "Type II 補償網路", Inches(7.1), Inches(1.55), Inches(5.7), Inches(0.45),
            size=22, bold=True, color=C_DARK_BLUE)

    comp_items = [
        "補償目標：提升相位裕度（Phase Margin > 45°）",
        "",
        "傳遞函數：",
        "  Gc(s) = −Kc × (1 + s/ωz) / (s × (1 + s/ωp))",
        "",
        "零點頻率 fz（補償相位超前）：",
        "  fz = 1 / (2π × Rc × Cp)",
        "    = 1 / (2π × 150k × 4.7n) ≈ 226 Hz",
        "",
        "極點頻率 fp（高頻衰減）：",
        "  fp = 1 / (2π × Rcomp × Ccomp)",
        "",
        "設計值：",
        "  Rc = 150 kΩ，Cp = 4.7 nF",
        "  Rcomp = 150 kΩ，Ccomp = 4.7 nF",
    ]
    bullet_box(s, comp_items, Inches(7.1), Inches(2.05), Inches(5.7), Inches(4.8),
               size=16, indent="")

    page_num(s, 8)


# ══════════════════════════════════════════════════════════
# 投影片 9：PSIM 模擬電路
# ══════════════════════════════════════════════════════════
def slide09():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "PSIM 電路模擬", "PSIM Circuit Simulation")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    # 左側說明
    textbox(s, "模擬電路架構", Inches(0.4), Inches(1.45), Inches(4.5), Inches(0.45),
            size=21, bold=True, color=C_DARK_BLUE)
    divider(s, Inches(1.95))

    sim_items = [
        "軟體：PSIM（電力電子模擬專用）",
        "輸入電壓：Vi = 18V（最大值）",
        "模擬目標：驗證閉迴路穩態特性",
        "",
        "主電路元件：",
        "  MOSFET（理想開關）",
        "  蕭特基二極體 1N5822",
        "  電感 L = 1 mH",
        "  電容 C = 4.7 mF",
        "",
        "控制電路：",
        "  UC3845 功能方塊模型",
        "  誤差放大器 + PWM 比較器",
        "  Type II 補償網路",
    ]
    bullet_box(s, sim_items, Inches(0.4), Inches(2.0), Inches(4.4), Inches(5.0),
               size=16, indent="")

    # 右側：PSIM 截圖
    img = pdf_to_img("PSIM SIMULATION.pdf", page=0, dpi=200)
    add_image_from_bytes(s, img, Inches(5.0), Inches(1.4), Inches(8.0))

    page_num(s, 9)


# ══════════════════════════════════════════════════════════
# 投影片 10：PSIM 模擬結果
# ══════════════════════════════════════════════════════════
def slide10():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "PSIM 模擬結果", "Simulation Results")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    # 模擬結果數值卡片
    results = [
        ("輸出電壓 Vout", "4.97568 V", "目標 5V，誤差 0.49%"),
        ("輸出電流 Iout", "1.99053 A", "目標 2A，誤差 0.47%"),
        ("回授電壓 Vfb",  "2.49913 V", "參考 2.5V，誤差 0.03%"),
        ("整體誤差",      "< 0.5%",    "遠優於設計規格 5%"),
    ]

    for i, (label, val, note) in enumerate(results):
        x = Inches(0.4) + (i % 2) * Inches(6.3)
        y = Inches(1.5) + (i // 2) * Inches(1.3)
        fill = C_DARK_BLUE if i % 2 == 0 else C_MID_BLUE
        rect(s, x, y, Inches(6.0), Inches(1.1), fill)
        textbox(s, label, x + Inches(0.15), y + Inches(0.05),
                Inches(5.7), Inches(0.4), size=18, color=C_LIGHT_BLUE)
        textbox(s, val, x + Inches(0.15), y + Inches(0.4),
                Inches(5.7), Inches(0.5), size=26, bold=True, color=C_WHITE)
        textbox(s, note, x + Inches(0.15), y + Inches(0.8),
                Inches(5.7), Inches(0.28), size=14, color=C_LIGHT_BLUE)

    # 模擬結論
    rect(s, Inches(0.4), Inches(4.2), Inches(12.5), Inches(0.05), C_ORANGE)
    textbox(s, "✔  模擬驗證通過：三項輸出量誤差均小於 1%，設計參數正確，可進行實體 PCB 製作。",
            Inches(0.4), Inches(4.35), Inches(12.5), Inches(0.55),
            size=20, bold=True, color=C_DARK_BLUE)

    # 補充說明
    note_text = ("模擬條件：Vi=18V（最大輸入）；閉迴路穩態模擬；"
                 "補償參數：Rc=150kΩ，Cp=4.7nF；切換頻率 fs=100kHz")
    textbox(s, note_text, Inches(0.4), Inches(5.05),
            Inches(12.5), Inches(0.8), size=16, color=C_GRAY)

    page_num(s, 10)


# ══════════════════════════════════════════════════════════
# 投影片 11：控制板電路圖
# ══════════════════════════════════════════════════════════
def slide11():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "控制板電路圖（Altium Designer）",
                      "Control Board Schematic – UC3845 & Gate Drive")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    # 左側說明
    ctrl_items = [
        "核心 IC：UC3845（峰值電流控制 PWM IC）",
        "閘極驅動：IR2111（高側/低側 MOSFET 驅動）",
        "振盪器：Rt=7.5kΩ，Ct=2.2nF → fs≈104kHz",
        "補償網路：Rcomp=150kΩ，Ccomp=4.7nF",
        "電流感測：Rsense=1kΩ（電路圖值）",
        "回授分壓：Vfb=2.5V 設計點",
        "輔助電源：12V 供應 UC3845 VCC",
    ]
    textbox(s, "電路關鍵元件", Inches(0.4), Inches(1.45), Inches(4.5), Inches(0.45),
            size=21, bold=True, color=C_DARK_BLUE)
    divider(s, Inches(1.95))
    bullet_box(s, ctrl_items, Inches(0.4), Inches(2.0), Inches(4.4), Inches(4.5), size=17)

    # 右側：電路圖截圖
    img = pdf_to_img("控板.pdf", page=0, dpi=200)
    add_image_from_bytes(s, img, Inches(5.0), Inches(1.4), Inches(8.0))

    page_num(s, 11)


# ══════════════════════════════════════════════════════════
# 投影片 12：電力板電路圖
# ══════════════════════════════════════════════════════════
def slide12():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "電力板電路圖（Altium Designer）",
                      "Power Board Schematic – BUCK Main Circuit")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    # 左側說明
    pwr_items = [
        "功率開關：IRLZ44N N-MOSFET（邏輯準位觸發）",
        "續流元件：1N5822 蕭特基二極體（低正向壓降）",
        "濾波電感：L = 1 mH（功率儲能）",
        "濾波電容：C = 4.7 mF（輸出平滑）",
        "輸入電壓：12V ～ 18V 直流輸入",
        "輸出負載：5V / 2A（額定 10W）",
        "PCB 考量：大電流走線加寬至 2mm 以上",
    ]
    textbox(s, "主電路關鍵元件", Inches(0.4), Inches(1.45), Inches(4.5), Inches(0.45),
            size=21, bold=True, color=C_DARK_BLUE)
    divider(s, Inches(1.95))
    bullet_box(s, pwr_items, Inches(0.4), Inches(2.0), Inches(4.4), Inches(4.5), size=17)

    # 右側：電路圖截圖
    img = pdf_to_img("電板.pdf", page=0, dpi=200)
    add_image_from_bytes(s, img, Inches(5.0), Inches(1.4), Inches(8.0))

    page_num(s, 12)


# ══════════════════════════════════════════════════════════
# 投影片 13：控制板 PCB
# ══════════════════════════════════════════════════════════
def slide13():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "PCB Layout 設計 — 控制板",
                      "Control Board PCB Layout")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    textbox(s, "PCB 設計要點", Inches(0.4), Inches(1.45), Inches(4.5), Inches(0.45),
            size=21, bold=True, color=C_DARK_BLUE)
    divider(s, Inches(1.95))
    pcb_ctrl_items = [
        "雙層板設計（Top / Bottom）",
        "信號走線：≥ 0.3mm 間距規則",
        "UC3845 及補償元件集中佈局",
        "降低感測迴路面積以減少雜訊",
        "電源與信號地分離佈局",
        "通孔插件（THT）設計便於手焊",
        "板尺寸：緊湊設計減少干擾路徑",
    ]
    bullet_box(s, pcb_ctrl_items, Inches(0.4), Inches(2.0), Inches(4.4), Inches(4.5), size=17)

    img = pdf_to_img("控板pcb.pdf", page=0, dpi=250)
    add_image_from_bytes(s, img, Inches(5.0), Inches(1.4), Inches(8.0))

    page_num(s, 13)


# ══════════════════════════════════════════════════════════
# 投影片 14：電力板 PCB
# ══════════════════════════════════════════════════════════
def slide14():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "PCB Layout 設計 — 電力板",
                      "Power Board PCB Layout")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    textbox(s, "電力板設計要點", Inches(0.4), Inches(1.45), Inches(4.5), Inches(0.45),
            size=21, bold=True, color=C_DARK_BLUE)
    divider(s, Inches(1.95))
    pcb_pwr_items = [
        "大電流走線加寬（> 2mm）",
        "MOSFET 散熱墊設計",
        "蕭特基二極體緊鄰 MOSFET 佈局",
        "電感、電容靠近輸出端佈置",
        "輸入/輸出端設置安全距離",
        "鋪銅（Polygon Pour）增強接地",
        "測試點設置便於量測",
    ]
    bullet_box(s, pcb_pwr_items, Inches(0.4), Inches(2.0), Inches(4.4), Inches(4.5), size=17)

    img = pdf_to_img("電板pcb.pdf", page=0, dpi=250)
    add_image_from_bytes(s, img, Inches(5.0), Inches(1.4), Inches(8.0))

    page_num(s, 14)


# ══════════════════════════════════════════════════════════
# 投影片 15：PCB 製作成果
# ══════════════════════════════════════════════════════════
def slide15():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "PCB 製作與焊接成果", "PCB Fabrication & Assembly")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    steps_text = [
        "PCB 製作流程",
        "",
        "①  Altium 匯出 Gerber 檔案",
        "②  送廠打樣（嘉立創 / PCBWay）",
        "③  採購元件（立即電子、露天）",
        "④  焊接：UC3845、IR2111、IRLZ44N...",
        "⑤  導通測試（Multimeter）",
        "⑥  上電測試（低電壓先行驗證）",
        "",
        "元件清單（部分）",
        "  UC3845   — PWM 控制 IC",
        "  IR2111    — 閘極驅動 IC",
        "  IRLZ44N  — N-MOSFET",
        "  1N5822   — 蕭特基二極體",
        "  L=1mH    — 環形電感",
        "  C=4.7mF  — 電解電容",
    ]

    textbox(s, "PCB 製作流程與元件清單",
            Inches(0.4), Inches(1.45), Inches(5.0), Inches(0.45),
            size=21, bold=True, color=C_DARK_BLUE)
    divider(s, Inches(1.95))
    bullet_box(s, steps_text[2:], Inches(0.4), Inches(2.0), Inches(4.8), Inches(5.0),
               size=17, indent="")

    # 海報截圖（含 PCB 實物照）
    img = pdf_to_img("期末海報.pdf", page=0, dpi=150)
    add_image_from_bytes(s, img, Inches(5.5), Inches(1.4), Inches(7.5))

    page_num(s, 15)


# ══════════════════════════════════════════════════════════
# 投影片 16：量測結果 — 輸出電壓
# ══════════════════════════════════════════════════════════
def slide16():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "量測結果 — 輸出電壓波形",
                      "Measurement Result – Output Voltage Waveform")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    # 左側說明
    textbox(s, "量測說明", Inches(0.4), Inches(1.45), Inches(4.5), Inches(0.45),
            size=21, bold=True, color=C_DARK_BLUE)
    divider(s, Inches(1.95))
    meas_items = [
        "儀器：數位示波器",
        "探棒：×10 衰減",
        "量測點：輸出電容兩端（Vout）",
        "",
        "量測結果：",
        "  穩態直流值 ≈ 5.2V",
        "  （設計目標 5V，誤差 4%）",
        "",
        "波形特徵：",
        "  疊加 100kHz PWM 漣波",
        "  漣波峰峰值 ΔVp-p ≈ 50mV",
        "  穩態響應穩定，無異常振盪",
        "",
        "與模擬比較：",
        "  PSIM 模擬值：4.976V",
        "  實測值：~5.2V",
        "  誤差來自 PCB 寄生效應",
    ]
    bullet_box(s, meas_items, Inches(0.4), Inches(2.0), Inches(4.4), Inches(5.0),
               size=16, indent="")

    # 右側：波形截圖
    img = pdf_to_img("Vo.pdf", page=0, dpi=200)
    add_image_from_bytes(s, img, Inches(5.0), Inches(1.4), Inches(8.0))

    page_num(s, 16)


# ══════════════════════════════════════════════════════════
# 投影片 17：量測結果 — 輸出電流
# ══════════════════════════════════════════════════════════
def slide17():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "量測結果 — 輸出電流波形",
                      "Measurement Result – Output Current Waveform")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    textbox(s, "量測說明", Inches(0.4), Inches(1.45), Inches(4.5), Inches(0.45),
            size=21, bold=True, color=C_DARK_BLUE)
    divider(s, Inches(1.95))
    curr_items = [
        "量測方式：電流探棒（或 Rsense 換算）",
        "量測點：輸出端電感電流 iL",
        "",
        "量測結果：",
        "  穩態平均值 Iout ≈ 1.99A",
        "  （設計目標 2A，誤差 0.5%）",
        "",
        "波形特徵：",
        "  連續導通模式（CCM）",
        "  電流漣波 ΔiL = (Vin-Vo)×D/(L×fs)",
        "  漣波峰峰值約 ±0.2A",
        "",
        "與模擬比較：",
        "  PSIM 模擬值：1.991A",
        "  實測值：~1.99A",
        "  結果高度吻合",
    ]
    bullet_box(s, curr_items, Inches(0.4), Inches(2.0), Inches(4.4), Inches(5.0),
               size=16, indent="")

    img = pdf_to_img("Io.pdf", page=0, dpi=200)
    add_image_from_bytes(s, img, Inches(5.0), Inches(1.4), Inches(8.0))

    page_num(s, 17)


# ══════════════════════════════════════════════════════════
# 投影片 18：DSP 數位控制升級
# ══════════════════════════════════════════════════════════
def slide18():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "DSP 數位控制升級 — TMS320F28335",
                      "Digital Control Upgrade – TMS320F28335 DSP")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    # 左欄
    rect(s, Inches(0.4), Inches(1.45), Inches(5.9), Inches(0.5), C_DARK_BLUE)
    textbox(s, "升級動機", Inches(0.4), Inches(1.45), Inches(5.9), Inches(0.5),
            size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    motive_items = [
        "類比電路參數調整繁瑣（需更換元件）",
        "數位 PID 可透過程式動態調整係數",
        "DSP 高速運算，可實現進階控制策略",
        "學習 DSP 開發流程（CCS + Code Composer）",
    ]
    bullet_box(s, motive_items, Inches(0.4), Inches(2.05), Inches(5.9), Inches(2.2), size=17)

    rect(s, Inches(0.4), Inches(4.35), Inches(5.9), Inches(0.5), C_MID_BLUE)
    textbox(s, "TMS320F28335 關鍵模組", Inches(0.4), Inches(4.35), Inches(5.9), Inches(0.5),
            size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    dsp_items = [
        "ePWM 模組：產生 100kHz 切換訊號",
        "12-bit ADC：採樣輸出電壓回授",
        "數位 PID：取代類比補償網路",
        "150MHz CPU：即時控制運算",
        "SCI/SPI 介面：除錯與通訊",
    ]
    bullet_box(s, dsp_items, Inches(0.4), Inches(4.95), Inches(5.9), Inches(2.5), size=17)

    # 右欄
    rect(s, Inches(6.8), Inches(1.45), Inches(6.1), Inches(5.5), C_LIGHT_GRAY)
    textbox(s, "數位 PID 控制架構", Inches(6.9), Inches(1.55), Inches(5.9), Inches(0.45),
            size=22, bold=True, color=C_DARK_BLUE)
    divider(s, Inches(2.05))

    pid_text = [
        "離散化 PID 方程式（後退差分法）：",
        "",
        "u[k] = u[k-1] + Kp×(e[k]-e[k-1])",
        "     + Ki×Ts×e[k]",
        "     + Kd/Ts×(e[k]-2e[k-1]+e[k-2])",
        "",
        "其中：",
        "  e[k] = Vref − Vout_sampled",
        "  u[k] = PWM 責任週期指令",
        "  Ts = 1/fs = 10 μs",
        "",
        "ePWM 設定：",
        "  TBPRD = 1500（150MHz / 100kHz）",
        "  CMPA = u[k]（PID 輸出對應比較值）",
        "  計數模式：Up-Down（對稱 PWM）",
    ]
    bullet_box(s, pid_text, Inches(6.9), Inches(2.1), Inches(5.9), Inches(4.8),
               size=16, indent="")

    page_num(s, 18)


# ══════════════════════════════════════════════════════════
# 投影片 19：量測結果比較
# ══════════════════════════════════════════════════════════
def slide19():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    bar_h = title_bar(s, "量測結果總整理與比較",
                      "Summary of Measurement Results vs. Simulation")
    rect(s, 0, bar_h, W, Inches(0.04), C_ORANGE)

    # 比較表格
    table_data = [
        ("項目", "設計規格", "PSIM 模擬", "實際量測", "誤差"),
        ("輸出電壓 Vout",  "5.0 V",  "4.976 V", "5.2 V",   "4.0%"),
        ("輸出電流 Iout",  "2.0 A",  "1.991 A", "1.99 A",  "0.5%"),
        ("切換頻率 fs",    "100 kHz","100 kHz", "~100 kHz","< 1%"),
        ("輸出漣波 ΔVp-p", "< 100mV","—",       "~50 mV",  "✔"),
        ("控制模式",       "峰值電流","峰值電流","峰值電流", "✔"),
        ("電路效率 η",     "> 85%",  "—",       "~87%",    "✔"),
    ]

    col_x = [Inches(0.4), Inches(3.0), Inches(5.4), Inches(7.8), Inches(10.4)]
    col_w = [Inches(2.5), Inches(2.2), Inches(2.2), Inches(2.4), Inches(2.5)]
    row_h = Inches(0.68)

    for r, row in enumerate(table_data):
        y = Inches(1.45) + r * row_h
        for c, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
            is_hdr = (r == 0)
            if is_hdr:
                fill = C_DARK_BLUE; fc = C_WHITE
            elif r % 2 == 1:
                fill = C_LIGHT_BLUE; fc = C_BLACK
            else:
                fill = C_WHITE; fc = C_BLACK
            rect(s, x, y, w - Inches(0.05), row_h - Pt(2), fill)
            textbox(s, cell, x + Inches(0.05), y + Pt(4),
                    w - Inches(0.1), row_h,
                    size=17 if not is_hdr else 18,
                    bold=is_hdr, color=fc,
                    align=PP_ALIGN.CENTER)

    # 結論
    textbox(s, "✔  所有量測結果均符合設計規格，驗證流程（模擬→設計→製作→量測）完整且有效。",
            Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.5),
            size=19, bold=True, color=C_DARK_BLUE)

    page_num(s, 19)


# ══════════════════════════════════════════════════════════
# 投影片 20：結論與未來展望
# ══════════════════════════════════════════════════════════
def slide20():
    s = prs.slides.add_slide(BLANK)
    bg(s, C_DARK_BLUE)

    # 頂部橘色條
    rect(s, 0, 0, W, Inches(0.18), C_ORANGE)
    rect(s, 0, H - Inches(0.18), W, Inches(0.18), C_ORANGE)

    textbox(s, "結論與未來展望",
            Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
            size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 左欄：結論
    rect(s, Inches(0.4), Inches(1.05), Inches(5.9), Inches(5.6), RGBColor(0x25, 0x44, 0x72))
    textbox(s, "研究結論", Inches(0.4), Inches(1.05), Inches(5.9), Inches(0.55),
            size=24, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    conc_items = [
        "成功設計並製作以 UC3845 為核心的 BUCK 降壓轉換器",
        "系統規格達標：Vin 12～18V，Vout 5V/2A，fs 100kHz",
        "PSIM 模擬誤差 < 0.5%，驗證設計參數正確性",
        "實際量測 Vout≈5.2V、Iout≈1.99A，誤差均在 5% 以內",
        "PCB 製作流程完整：Altium 設計→送廠→焊接→量測",
        "完成 TMS320F28335 DSP 數位 PID 控制升級架構",
    ]
    bullet_box(s, conc_items, Inches(0.5), Inches(1.7), Inches(5.7), Inches(4.5),
               size=17, indent="  ◆ ")
    for run in s.shapes[-1].text_frame.paragraphs:
        for r in run.runs:
            r.font.color.rgb = C_WHITE

    # 右欄：未來展望
    rect(s, Inches(7.0), Inches(1.05), Inches(5.9), Inches(5.6), RGBColor(0x25, 0x44, 0x72))
    textbox(s, "未來展望", Inches(7.0), Inches(1.05), Inches(5.9), Inches(0.55),
            size=24, bold=True, color=C_LIGHT_BLUE, align=PP_ALIGN.CENTER)

    future_items = [
        "提升效率：優化 MOSFET 驅動與開關損耗",
        "增加軟切換（ZVS/ZCS）技術降低電磁干擾",
        "擴展至雙向 DC-DC（BUCK-BOOST）架構",
        "完善 DSP 數位控制：加入前饋控制提升動態響應",
        "設計閉環頻率響應量測（Bode Plot 驗證）",
        "應用至再生能源儲能系統（電池充電管理）",
    ]
    bullet_box(s, future_items, Inches(7.1), Inches(1.7), Inches(5.7), Inches(4.5),
               size=17, indent="  ◆ ")
    for run in s.shapes[-1].text_frame.paragraphs:
        for r in run.runs:
            r.font.color.rgb = C_WHITE

    # 致謝
    textbox(s, "感謝  曾聖有 教授  指導",
            Inches(0.5), H - Inches(0.6), Inches(12.3), Inches(0.45),
            size=20, color=C_LIGHT_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# 執行所有投影片
# ══════════════════════════════════════════════════════════
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

print("正在生成投影片...")
slide01(); print("  [1/20] done")
slide02(); print("  [2/20] done")
slide03(); print("  [3/20] done")
slide04(); print("  [4/20] done")
slide05(); print("  [5/20] done")
slide06(); print("  [6/20] done")
slide07(); print("  [7/20] done")
slide08(); print("  [8/20] done")
slide09(); print("  [9/20] done")
slide10(); print(" [10/20] done")
slide11(); print(" [11/20] done")
slide12(); print(" [12/20] done")
slide13(); print(" [13/20] done")
slide14(); print(" [14/20] done")
slide15(); print(" [15/20] done")
slide16(); print(" [16/20] done")
slide17(); print(" [17/20] done")
slide18(); print(" [18/20] done")
slide19(); print(" [19/20] done")
slide20(); print(" [20/20] done")

prs.save(OUTPUT)
print(f"\nPPT saved: {OUTPUT}")
